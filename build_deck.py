#!/usr/bin/env python3
"""Restructure Portal Engineering investor deck: secondary 25% sale + promoter loan.
Enhances existing template, adds IB-grade slides, fixes typos, adds caveats."""
import copy, zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = "/home/user/kulkarnivishal1/portal_deck.pptx"
OUT = "/home/user/kulkarnivishal1/PORTAL_ENGINEERING_Investor_Deck_v2.pptx"

# ---- palette ----
DARK   = "15092E"
WHITE  = "FFFFFF"
GOLD   = "FBC76E"
LPURP  = "C8B8FF"   # light purple subtitle
MPURP  = "3A2860"   # mid purple subtext
PANEL  = "241147"   # deep purple panel
CARD_L = "FFFFFF"   # light card
CARD_L2= "F0ECF8"
HEAD   = "Montserrat"
LIGHT  = "Montserrat Light"

SW, SH = 20.0, 11.25

prs = Presentation(SRC)
BLANK = prs.slides[0].slide_layout

# extract bg images
z = zipfile.ZipFile(SRC)
import os
os.makedirs("/tmp/media", exist_ok=True)
for nm in ["image1","image5"]:
    open(f"/tmp/media/{nm}.png","wb").write(z.read(f"ppt/media/{nm}.bin"))
BG_DARK = "/tmp/media/image1.png"
BG_GRAD = "/tmp/media/image5.png"

def new_slide():
    s = prs.slides.add_slide(BLANK)
    for ph in list(s.shapes):
        ph._element.getparent().remove(ph._element)
    s.shapes.add_picture(BG_DARK, 0, 0, Inches(SW), Inches(SH))
    s.shapes.add_picture(BG_GRAD, 0, 0, Inches(SW), Inches(SH))
    return s

def txt(s, l,t,w,h, text, size, color, bold=False, font=HEAD,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=None, sb=None):
    tb = s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in (tf.margin_left,):
        pass
    tf.margin_left=Inches(0.06); tf.margin_right=Inches(0.06)
    tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    for i,ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        if ls: p.line_spacing=ls
        if sb is not None: p.space_before=Pt(sb)
        r=p.add_run(); r.text=ln
        r.font.size=Pt(size); r.font.bold=bold; r.font.name=font
        r.font.color.rgb=RGBColor.from_string(color)
    return tb

def rect(s, l,t,w,h, fill, rounded=False, line=None, lw=1.0, adj=0.08):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb=RGBColor.from_string(fill)
    if line:
        shp.line.color.rgb=RGBColor.from_string(line); shp.line.width=Pt(lw)
    else:
        shp.line.fill.background()
    shp.shadow.inherit=False
    if rounded:
        try: shp.adjustments[0]=adj
        except Exception: pass
    return shp

def header(s, label, title, subtitle):
    rect(s, 0.65, 0.80, 0.07, 0.73, GOLD)
    txt(s, 1.33, 0.55, 17.5, 0.32, label, 15, GOLD, bold=True, font=HEAD)
    txt(s, 1.30, 0.86, 17.5, 0.78, title, 40, WHITE, bold=False, font=LIGHT)
    if subtitle:
        txt(s, 1.33, 1.74, 17.5, 0.50, subtitle, 18, LPURP, bold=False, font=LIGHT)

def caveat(s, y=10.62):
    txt(s, 1.33, y, 17.5, 0.34,
        "All projections are indicative and subject to availability of funds at the right time.",
        11.5, LPURP, bold=False, font=LIGHT, align=PP_ALIGN.LEFT)

# =====================================================================
# SLIDE A — KEY INVESTMENT HIGHLIGHTS
# =====================================================================
sA = new_slide()
header(sA, "WHY PORTAL ENGINEERING", "Key Investment Highlights",
       "A profitable, fast-growing import-substitution play in India's elevator & parking systems market")
tiles = [
    ("15+ Yrs","Proven Track Record","Established maker of automatic lift doors, hydraulic systems & multi-level parking"),
    ("~52%","Revenue CAGR","Y-o-Y revenue growth sustained over FY22–FY25"),
    ("100%+","PAT Growth","Profit after tax more than doubled over the last two years"),
    ("3 Geographies","Live Export Traction","Recurring orders from Kuwait, Canada & Australia"),
    ("New Plant","Capacity Scale-Up","Ambernath facility + advanced machinery driving margin gains"),
    ("25% Stake","Secondary on Offer","Promoter stake sale at ~₹95 Cr equity value"),
]
cw, ch, gx, gy = 5.78, 3.10, 0.40, 0.45
x0, y0 = 0.82, 2.55
for i,(num,lbl,desc) in enumerate(tiles):
    r,c = divmod(i,3)
    x = x0 + c*(cw+gx); y = y0 + r*(ch+gy)
    rect(sA, x, y, cw, ch, CARD_L, rounded=True, adj=0.06)
    rect(sA, x, y, 0.14, ch, GOLD)  # accent edge
    txt(sA, x+0.30, y+0.28, cw-0.5, 0.95, num, 38, DARK, bold=True, font=HEAD)
    txt(sA, x+0.30, y+1.28, cw-0.5, 0.45, lbl, 17, DARK, bold=True, font=HEAD)
    txt(sA, x+0.30, y+1.80, cw-0.55, 1.15, desc, 13.5, MPURP, bold=False, font=LIGHT, ls=1.05)
caveat(sA)

# =====================================================================
# SLIDE B — TRANSACTION STRUCTURE
# =====================================================================
sB = new_slide()
header(sB, "THE PROPOSAL", "Transaction Structure",
       "A secondary stake sale — funded straight back into the business as interest-free promoter capital")
steps = [
    ("1","Secondary Sale","PE investor acquires 25% equity by buying existing shares from the promoter. No new shares are issued — zero primary dilution."),
    ("2","Proceeds to Promoter","Promoter receives ₹23.75 Cr. The company's share capital and balance sheet are unchanged at this step."),
    ("3","Re-Infused as Loan","Promoter infuses the entire ₹23.75 Cr back into Portal as an interest-free, unsecured loan (quasi-equity)."),
    ("4","Repaid from Cash Flow","Company deploys the capital for growth; the loan is drawn down from internal free cash flows over 3–5 years."),
]
cw, ch = 4.05, 4.55
gap = 0.70
x0, y = 0.82, 3.55
for i,(n,h,b) in enumerate(steps):
    x = x0 + i*(cw+gap)
    rect(sB, x, y, cw, ch, CARD_L, rounded=True, adj=0.05)
    rect(sB, x, y, cw, 0.14, GOLD)
    # number badge
    bd = rect(sB, x+0.30, y+0.36, 0.78, 0.78, DARK, rounded=True, adj=0.5)
    txt(sB, x+0.30, y+0.46, 0.78, 0.60, n, 26, GOLD, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
    txt(sB, x+0.30, y+1.40, cw-0.6, 0.80, h, 19, DARK, bold=True, font=HEAD)
    txt(sB, x+0.30, y+2.30, cw-0.62, 2.0, b, 14, MPURP, bold=False, font=LIGHT, ls=1.08)
    if i < 3:
        ar = sB.shapes.add_shape(MSO_SHAPE.CHEVRON,
              Inches(x+cw+0.10), Inches(y+ch/2-0.30), Inches(0.50), Inches(0.60))
        ar.fill.solid(); ar.fill.fore_color.rgb=RGBColor.from_string(GOLD)
        ar.line.fill.background(); ar.shadow.inherit=False
# summary band
by = 8.55
rect(sB, 0.82, by, 18.36, 1.30, PANEL, rounded=True, adj=0.10)
band = [("PE acquires","25% stake"),("for consideration of","₹23.75 Cr"),
        ("Implied 100% equity value","~₹95 Cr"),("Capital available to business","₹23.75 Cr")]
bw = 18.36/4
for i,(a,b) in enumerate(band):
    bx = 0.82 + i*bw
    txt(sB, bx, by+0.20, bw, 0.40, a, 13, LPURP, bold=False, font=LIGHT, align=PP_ALIGN.CENTER)
    txt(sB, bx, by+0.58, bw, 0.55, b, 23, GOLD, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
    if i>0:
        rect(sB, bx-0.01, by+0.25, 0.012, 0.80, GOLD)
caveat(sB)

# =====================================================================
# SLIDE C — INDICATIVE TRANSACTION TERMS
# =====================================================================
sC = new_slide()
header(sC, "DEAL SNAPSHOT", "Indicative Transaction Terms",
       "Key commercial terms of the proposed 25% secondary")
rows = [
    ("Transaction Type","Secondary sale of existing shares — no primary dilution"),
    ("Stake on Offer","25%"),
    ("Consideration","₹23.75 Crores"),
    ("Implied Equity Value (100%)","~₹95 Crores"),
    ("Implied Valuation Multiple","~17.6x FY26-27E PAT  |  ~9.3x FY27-28E PAT"),
    ("Re-Investment Instrument","Interest-free, unsecured promoter loan (quasi-equity)"),
    ("Loan Repayment","Drawn from internal accruals over 3–5 years"),
    ("Use of Capital","Growth capex, capacity, inventory, certifications & branding"),
    ("Governance","1 board seat + customary information & minority rights (indicative)"),
]
rx, ry = 0.82, 2.70
rh, rgap = 0.74, 0.10
lw_, vw_ = 5.7, 12.66
for i,(a,b) in enumerate(rows):
    y = ry + i*(rh+rgap)
    rect(sC, rx, y, lw_, rh, PANEL, rounded=True, adj=0.12)
    rect(sC, rx+lw_+0.16, y, vw_, rh, CARD_L if i%2==0 else CARD_L2, rounded=True, adj=0.12)
    txt(sC, rx+0.30, y, lw_-0.5, rh, a, 15, LPURP, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
    txt(sC, rx+lw_+0.46, y, vw_-0.7, rh, b, 15, DARK, bold=False, font=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
caveat(sC, y=10.45)

# =====================================================================
# EDIT EXISTING SLIDES
# =====================================================================
def set_para_text(shape, new_text):
    """Set text of a text frame's first paragraph, keep first run's formatting, drop extra runs."""
    tf = shape.text_frame
    p = tf.paragraphs[0]
    runs = p.runs
    if not runs:
        r = p.add_run(); r.text=new_text; return
    runs[0].text = new_text
    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)

def find(slide, name):
    for sh in slide.shapes:
        if sh.name == name and sh.has_text_frame:
            return sh
    return None

# --- Slide 26 (index 25): Use of Funds reframe ---
s26 = prs.slides[25]
for sh in s26.shapes:
    if not sh.has_text_frame: continue
    t = sh.text_frame.text.strip()
    if t == "Ask and Use of Funds":
        set_para_text(sh, "Use of Funds")
    elif t.startswith("Allocation of INR 23.75"):
        set_para_text(sh, "Deployment of the ₹23.75 Cr promoter loan into growth capex & working capital")
    elif t == "FUNDS ALLOCATION BREAKDOWN":
        set_para_text(sh, "USE OF FUNDS  —  TOTAL ₹23.75 CRORES")

# --- Slide 31 (index 30): reframe immediate infusion as first tranche ---
s31 = prs.slides[30]
m = {
 "main-title":"Near-Term Impact: First Tranche",
 "sub-title":"What an immediate ₹5 Cr first tranche unlocks in H1 FY 26-27  |  Subject to timely infusion",
 "card3-lbl":"First Tranche",
 "card3-sub":"First draw of the ₹23.75 Cr promoter loan",
}
for sh in s31.shapes:
    if sh.has_text_frame and sh.name in m:
        set_para_text(sh, m[sh.name])
# fix pre-existing low-contrast: card3 had white text on a gold card -> make dark
c3 = {"card3-num":DARK, "card3-lbl":DARK, "card3-sub":MPURP, "bullets":DARK}
for sh in s31.shapes:
    if sh.has_text_frame and sh.name in c3:
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                r.font.color.rgb = RGBColor.from_string(c3[sh.name])

# --- Typos ---
# Slide 27 (index 26): EBITA -> EBITDA
s27 = prs.slides[26]
for sh in s27.shapes:
    if sh.has_text_frame and "EBITA" in sh.text_frame.text and "EBITDA" not in sh.text_frame.text.replace("EBITDA",""):
        if sh.text_frame.text.strip().startswith("Consistent"):
            set_para_text(sh, "Consistent 78%+ growth in EBITDA Y-o-Y over the past 3 years")
# Slide 33 (index 32): PRODUTIVITY -> PRODUCTIVITY
s33 = prs.slides[32]
for sh in s33.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip()=="PRODUTIVITY":
        set_para_text(sh, "PRODUCTIVITY")

# --- Caveat footnotes on projection slides 29, 30 ---
for idx in [28, 29]:  # slide 29 & 30
    caveat(prs.slides[idx], y=10.78)

# =====================================================================
# REORDER SLIDES
# new slides currently appended at indices 37(A),38(B),39(C)
# =====================================================================
sldIdLst = prs.slides._sldIdLst
els = list(sldIdLst)
# A=highlights ->after agenda(idx1); B,C ->before Use of Funds(idx25)
order = [0,1,37] + list(range(2,25)) + [38,39] + list(range(25,37))
assert sorted(order)==list(range(40)), sorted(order)
for i in order:
    sldIdLst.append(els[i])

prs.save(OUT)
print("SAVED", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
print("final count:", len(list(Presentation(OUT).slides)))
