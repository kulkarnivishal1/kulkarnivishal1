#!/usr/bin/env python3
"""Approximate PNG preview of pptx slides (layout verification; final deck uses Montserrat)."""
import sys
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageDraw, ImageFont

DECK="/home/user/kulkarnivishal1/PORTAL_ENGINEERING_Investor_Deck_v2.pptx"
SCALE=80  # px per inch
SW,SH=int(20*SCALE),int(11.25*SCALE)
REG="/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
BLD="/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
_fc={}
def font(sz,bold):
    k=(int(sz),bold)
    if k not in _fc: _fc[k]=ImageFont.truetype(BLD if bold else REG,int(sz))
    return _fc[k]

def emu_in(v): return Emu(v).inches if v is not None else 0

def runs_info(tf):
    paras=[]
    for para in tf.paragraphs:
        txt="".join(r.text for r in para.runs)
        r0=para.runs[0] if para.runs else None
        sz=r0.font.size.pt if (r0 and r0.font.size) else 14
        bold=bool(r0.font.bold) if r0 else False
        try: col=str(r0.font.color.rgb) if r0 else "FFFFFF"
        except: col="FFFFFF"
        al=para.alignment
        paras.append((txt,sz,bold,col,al))
    return paras

def wrap(draw,text,fnt,maxw):
    words=text.split(" "); lines=[]; cur=""
    for w in words:
        t=(cur+" "+w).strip()
        if draw.textlength(t,font=fnt)<=maxw or not cur: cur=t
        else: lines.append(cur); cur=w
    if cur: lines.append(cur)
    return lines

def render(idx,out):
    prs=Presentation(DECK); s=prs.slides[idx]
    img=Image.new("RGB",(SW,SH),(21,9,46)); d=ImageDraw.Draw(img,"RGBA")
    for sh in s.shapes:
        l=emu_in(sh.left)*SCALE; t=emu_in(sh.top)*SCALE
        w=emu_in(sh.width)*SCALE; h=emu_in(sh.height)*SCALE
        st=str(sh.shape_type)
        # background pictures
        if "PICTURE" in st:
            try:
                im=sh.image; from io import BytesIO
                pic=Image.open(BytesIO(im.blob)).convert("RGBA")
                pic=pic.resize((max(1,int(w)),max(1,int(h))))
                img.paste(pic,(int(l),int(t)),pic)
            except Exception: pass
            continue
        # autoshapes (rect/rounded/chevron) with solid fill
        if "AUTO_SHAPE" in st or "ROUNDED" in st:
            fill=None
            try:
                if sh.fill.type is not None:
                    fill=str(sh.fill.fore_color.rgb)
            except Exception: fill=None
            if fill:
                c=tuple(int(fill[i:i+2],16) for i in (0,2,4))
                d.rounded_rectangle([l,t,l+w,t+h],radius=min(14,h/3),fill=c)
        # text
        if sh.has_text_frame and sh.text_frame.text.strip():
            paras=runs_info(sh.text_frame)
            anchor=sh.text_frame.vertical_anchor
            # build wrapped lines
            alllines=[]
            for (txt,sz,bold,col,al) in paras:
                fnt=font(sz*SCALE/72.0,bold)
                for ln in wrap(d,txt,fnt,w-8):
                    alllines.append((ln,fnt,col,al,sz))
            lh=sum(f.getbbox("Ag")[3]-f.getbbox("Ag")[1]+f.size*0.35 for _,f,_,_,_ in alllines) if alllines else 0
            if anchor==MSO_ANCHOR.MIDDLE: cy=t+(h-lh)/2
            elif anchor==MSO_ANCHOR.BOTTOM: cy=t+h-lh
            else: cy=t+4
            for ln,fnt,col,al,sz in alllines:
                try: c=tuple(int(col[i:i+2],16) for i in (0,2,4))
                except: c=(255,255,255)
                tw=d.textlength(ln,font=fnt)
                if al==PP_ALIGN.CENTER: cx=l+(w-tw)/2
                elif al==PP_ALIGN.RIGHT: cx=l+w-tw-4
                else: cx=l+6
                d.text((cx,cy),ln,font=fnt,fill=c)
                cy+=fnt.getbbox("Ag")[3]-fnt.getbbox("Ag")[1]+fnt.size*0.35
    img.save(out); print("wrote",out)

if __name__=="__main__":
    for a in sys.argv[1:]:
        i=int(a); render(i,f"/tmp/prev_{i+1}.png")
